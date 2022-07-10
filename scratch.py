# OS module facilitates Operating System functionality like going through directory
import os
# SYS module allows tasks like exiting is performed using this module
import sys
# askdirectory module will be used to create folder choosing dialog boxes
from tkinter.filedialog import askdirectory
import collections  # pptx uses this module internally
import collections.abc  # pptx uses this module internally
# pptx used to create Presentation files
from pptx import Presentation
# datetime used to generate current date and time
from datetime import datetime


# Returns the current system time
def currentTime():
    return datetime.now().strftime("%H-%M-%S")


# Returns the current date from the system
def currentDate():
    return datetime.today().strftime('%Y-%m-%d')


# exitter() will take user prompt to exit out of python script
def exitter():
    quitter = input("Press ENTER to quit")
    if quitter:
        sys.exit(1)


ORIGIN_PATH = askdirectory(title='Select your original folder')  # Shows dialog box and return the path
if ORIGIN_PATH == "":
    print('You did not select the origin folder.\n'
          'Please quit and start the application, and do as instructed\n')
    exitter()

Destination_PATH = askdirectory(title='Please select the destination folder')
if Destination_PATH == "":
    print('You did not select the destination folder.\n'
          'Please quit and start the application, and do as instructed\n')
    exitter()

# Generate a filename based on current time and date
origin = ORIGIN_PATH.split('/')
filename = currentTime() + '-' + currentDate() + '.pptx'
for item in origin:
    # Changed the filename based on origin folder if the origin is valid
    if item.startswith('INT'):
        filename = item + ".pptx"

savingPATH = os.path.join(Destination_PATH, filename)

# Path to the folder that should contain "Output 20xxxxxx" folder
OutputFolder = "Output 20xxxxxx"
# Image name that needs to be placed to the PowerPoint
imageName = "pre_6_l.png"

# The folder that contains the images
imageFolder = "images"

# Creating the Presentation object and setting the layout
myPresentation = Presentation()
layout = myPresentation.slide_layouts[8]

# Creating a path to the folder that contains "Scan" folders
myPath = os.path.join(ORIGIN_PATH, OutputFolder)

# If the "Output 20xxxxxx" folder exists, list all the folders inside it
if os.path.exists(myPath):
    directories = os.listdir(myPath)
else:
    print('Please select a directory that contains the folder "Output 20xxxxxx folder')
    exitter()

# Slide counter
slide_count = 0
# Go through each scan folders and put the picture to one slide each in the PowerPoint file
for directory in directories:
    imagePath = os.path.join(myPath, directory, imageFolder, imageName)
    if os.path.exists(imagePath):
        slide = myPresentation.slides.add_slide(layout)
        slide_count += 1
        slide.shapes.title.text = f'Image from {directory}'
        slide.placeholders[1].insert_picture(imagePath)

# Saving the PowerPoint file and saving it
if slide_count > 0:
    print(f"{slide_count} image(s) added to the powerpoint presentation")
    myPresentation.save(savingPATH)

    # Opening the PowerPoint file
    os.startfile(savingPATH)
else:
    print('No images to add')
